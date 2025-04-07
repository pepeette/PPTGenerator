import os
import requests
import json
from pptx import Presentation
from pptx.util import Inches, Pt
import nltk
from nltk.tokenize import sent_tokenize, word_tokenize
from nltk.corpus import stopwords
import ssl
import streamlit as st
import tempfile

# Download NLTK resources (if not already downloaded)
try:
    _create_unverified_https_context = ssl._create_unverified_context
except AttributeError:
    pass
else:
    ssl._create_default_https_context = _create_unverified_https_context

nltk.download('punkt', quiet=True)
nltk.download('stopwords', quiet=True)

class ArticleToPPT:
    def __init__(self, api_key=None, api_endpoint="https://api.perplexity.ai/chat/completions"):
        self.api_key = api_key
        self.api_endpoint = api_endpoint
        self.output_path = "/Desktop/powerpoint"
        
    def set_api_key(self, api_key):
        self.api_key = api_key
        
    def extract_key_points(self, text):
        """Extract key points from text using simple NLP techniques"""
        sentences = sent_tokenize(text)
        
        # Calculate sentence importance based on word frequency
        words = word_tokenize(text.lower())
        stop_words = set(stopwords.words('english'))
        words = [word for word in words if word.isalnum() and word not in stop_words]
        
        word_freq = {}
        for word in words:
            if word in word_freq:
                word_freq[word] += 1
            else:
                word_freq[word] = 1
                
        # Calculate sentence scores
        sent_scores = {}
        for i, sentence in enumerate(sentences):
            sent_words = word_tokenize(sentence.lower())
            for word in sent_words:
                if word in word_freq:
                    if i in sent_scores:
                        sent_scores[i] += word_freq[word]
                    else:
                        sent_scores[i] = word_freq[word]
        
        # Get top sentences
        top_sentences = sorted(sent_scores.items(), key=lambda x: x[1], reverse=True)
        
        # Extract top 5-7 key points
        key_points = []
        for i, score in top_sentences[:min(7, len(top_sentences))]:
            key_points.append(sentences[i])
            
        return key_points
    
    def get_insights_from_llm(self, text):
        """Fetch insights from an LLM API"""
        if not self.api_key:
            st.warning("No API key provided for LLM insights. Using local extraction only.")
            return {
                "summary": "No LLM summary available - add your API key to enable this feature.",
                "key_points": self.extract_key_points(text),
                "title_suggestion": "Article Analysis"
            }
        
        try:
            headers = {
                "Authorization": f"Bearer {self.api_key}",
                "Content-Type": "application/json"
            }
            
            payload = {
                "model": "mixtral-8x7b-instruct",
                "messages": [
                    {
                        "role": "system", 
                        "content": "You are an expert at analyzing text and extracting key insights."
                    },
                    {
                        "role": "user", 
                        "content": f"Please analyze this text and provide: 1) A concise summary (3-4 sentences), 2) 5-7 key points formatted as bullet points, 3) A compelling title suggestion, 4) 3 discussion questions for students, 5) 2-3 categories or themes present in the text. Format your response as JSON with keys: summary, key_points (array), title_suggestion, discussion_questions (array), and themes (array). Here's the text: {text[:4000]}..."
                    }
                ],
                "max_tokens": 1000
            }
            
            with st.spinner("Fetching insights from LLM..."):
                response = requests.post(self.api_endpoint, headers=headers, json=payload)
                response_data = response.json()
            
            # Extract the response content from the LLM
            content = response_data.get("choices", [{}])[0].get("message", {}).get("content", "")
            
            # Parse the JSON response
            try:
                # Extract the JSON part from the response
                json_start = content.find('{')
                json_end = content.rfind('}') + 1
                json_str = content[json_start:json_end]
                
                insights = json.loads(json_str)
                return insights
            except json.JSONDecodeError:
                st.warning("Failed to parse LLM response as JSON. Using local extraction instead.")
                return {
                    "summary": content[:200] + "...",
                    "key_points": self.extract_key_points(text),
                    "title_suggestion": "Article Analysis",
                    "discussion_questions": ["What are the main ideas in this text?", 
                                            "How does this relate to real-world applications?",
                                            "What perspectives might be missing from this text?"],
                    "themes": ["Main Theme", "Secondary Theme"]
                }
        except Exception as e:
            st.error(f"Error fetching insights from LLM: {e}")
            return {
                "summary": "Error fetching LLM insights. See console for details.",
                "key_points": self.extract_key_points(text),
                "title_suggestion": "Article Analysis",
                "discussion_questions": ["What are the main ideas in this text?"],
                "themes": ["Main Theme"]
            }
    
    def create_presentation(self, text, insights):
        """Create a PowerPoint presentation from text and insights"""
        prs = Presentation()
        
        # Add title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = insights.get("title_suggestion", "Article Analysis")
        subtitle.text = "Insights and Key Points"
        
        # Add summary slide
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "Summary"
        content.text = insights.get("summary", "No summary available.")
        
        # Add key points slides
        key_points = insights.get("key_points", [])
        if not key_points:
            key_points = self.extract_key_points(text)
        
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "Key Points"
        content.text = ""
        for point in key_points:
            p = content.text_frame.add_paragraph()
            p.text = "• " + point
            p.level = 0
            
        # Add themes slide if available
        themes = insights.get("themes", [])
        if themes:
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            content = slide.placeholders[1]
            
            title.text = "Key Themes"
            content.text = ""
            for theme in themes:
                p = content.text_frame.add_paragraph()
                p.text = "• " + theme
                p.level = 0
                
        # Add discussion questions slide if available
        questions = insights.get("discussion_questions", [])
        if questions:
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            content = slide.placeholders[1]
            
            title.text = "Discussion Questions"
            content.text = ""
            for question in questions:
                p = content.text_frame.add_paragraph()
                p.text = "• " + question
                p.level = 0
        
        # Add final slide
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = "Thank You"
        
        # Save the presentation to a temp file first (for streamlit download button)
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp_file:
            prs.save(tmp_file.name)
            temp_path = tmp_file.name
        
        # Also save to the specified output directory
        home_dir = os.path.expanduser("~")
        output_dir = os.path.join(home_dir, self.output_path.lstrip("/"))
        os.makedirs(output_dir, exist_ok=True)
        
        # Generate a filename based on the title
        safe_title = ''.join(c if c.isalnum() or c in [' ', '-', '_'] else '_' for c in insights.get("title_suggestion", "Article_Analysis"))
        safe_title = safe_title.replace(' ', '_')
        filename = f"{safe_title}_{len(os.listdir(output_dir))}.pptx"
        filepath = os.path.join(output_dir, filename)
        
        # Copy from temp file to final destination
        import shutil
        shutil.copy2(temp_path, filepath)
        
        return filepath, temp_path, filename
    
    def process_text(self, text, use_llm=True):
        """Process text and generate a PowerPoint presentation"""
        if use_llm and self.api_key:
            insights = self.get_insights_from_llm(text)
        else:
            insights = {
                "summary": text[:200] + "...",
                "key_points": self.extract_key_points(text),
                "title_suggestion": "Article Analysis",
                "discussion_questions": ["What are the main ideas in this text?", 
                                        "How does this relate to real-world applications?",
                                        "What perspectives might be missing from this text?"],
                "themes": ["Main Theme", "Secondary Theme"]
            }
            
        filepath, temp_path, filename = self.create_presentation(text, insights)
        return filepath, temp_path, filename, insights

# Streamlit App
def main():
    st.set_page_config(page_title="Article to PowerPoint Generator", layout="wide")
    
    st.title("Article to PowerPoint Generator")
    st.markdown("Convert articles and text into professional PowerPoint presentations")
    
    # Initialize converter
    converter = ArticleToPPT()
    
    # Sidebar for configuration
    with st.sidebar:
        st.header("Configuration")
        api_key = st.text_input("API Key (optional)", type="password", 
                               help="Enter your Perplexity AI API key for better insights")
        
        use_llm = st.checkbox("Use LLM for insights", value=True, 
                             help="If checked and API key provided, the app will use an LLM to generate better insights")
        
        output_path = st.text_input("Output Path", value="/Desktop/powerpoint", 
                                    help="Path where presentations will be saved (relative to your home directory)")
        converter.output_path = output_path
        
        st.divider()
        st.markdown("### About")
        st.markdown("This app helps you quickly convert articles into well-structured PowerPoint presentations.")
        st.markdown("Perfect for educators and counselors who need to create professional material for students.")
    
    # Set API key if provided
    if api_key:
        converter.set_api_key(api_key)
    
    # Text input area
    st.subheader("Enter or paste your article")
    text_input_method = st.radio("Input method:", ["Type/Paste Text", "Upload File"])
    
    if text_input_method == "Type/Paste Text":
        text = st.text_area("", height=300, placeholder="Paste your article text here...")
    else:
        uploaded_file = st.file_uploader("Choose a text file", type=["txt", "md", "pdf"])
        text = ""
        if uploaded_file is not None:
            # Handle PDF files
            if uploaded_file.name.endswith('.pdf'):
                try:
                    import PyPDF2
                    pdf_reader = PyPDF2.PdfReader(uploaded_file)
                    text = ""
                    for page in pdf_reader.pages:
                        text += page.extract_text() + "\n"
                except ImportError:
                    st.error("PDF processing requires PyPDF2. Install with: pip install PyPDF2")
                    st.stop()
            else:
                # Handle text files
                text = uploaded_file.getvalue().decode("utf-8")
            
            st.success(f"File loaded: {uploaded_file.name}")
            st.text_area("Preview", text[:500] + "...", height=150)
    
    # Process button
    if st.button("Generate PowerPoint", disabled=not text):
        if not text.strip():
            st.error("Please enter or upload some text first.")
        else:
            with st.spinner("Processing... This may take a minute."):
                filepath, temp_path, filename, insights = converter.process_text(text, use_llm and api_key)
            
            # Display success message and insights
            st.success(f"PowerPoint presentation created successfully!")
            st.info(f"Saved to: {filepath}")
            
            # Download button
            with open(temp_path, "rb") as file:
                btn = st.download_button(
                    label="Download Presentation",
                    data=file,
                    file_name=filename,
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                )
            
            # Display insights
            with st.expander("View Generated Insights", expanded=True):
                st.subheader(insights.get("title_suggestion", "Article Analysis"))
                
                st.markdown("#### Summary")
                st.write(insights.get("summary", "No summary available."))
                
                st.markdown("#### Key Points")
                for point in insights.get("key_points", []):
                    st.markdown(f"• {point}")
                
                if "themes" in insights:
                    st.markdown("#### Key Themes")
                    for theme in insights.get("themes", []):
                        st.markdown(f"• {theme}")
                
                if "discussion_questions" in insights:
                    st.markdown("#### Discussion Questions")
                    for question in insights.get("discussion_questions", []):
                        st.markdown(f"• {question}")

if __name__ == "__main__":
    main()
